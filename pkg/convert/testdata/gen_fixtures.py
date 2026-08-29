#!/usr/bin/env python3
"""Generate deterministic test fixtures for pkg/convert golden tests.

Run: python3 gen_fixtures.py
Output: deterministic samples in this directory (git-committed).
"""
import os

HERE = os.path.dirname(os.path.abspath(__file__))

# --- PDF (reportlab) ---
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter


def gen_pdf():
    path = os.path.join(HERE, "sample.pdf")
    c = canvas.Canvas(path, pagesize=letter)
    c.setFont("Helvetica-Bold", 20)
    c.drawString(72, 720, "OKF PDF Fixture")
    c.setFont("Helvetica", 12)
    c.drawString(72, 690, "This is the body text of the PDF fixture.")
    c.drawString(72, 672, "It has two lines of plain content.")
    c.showPage()
    c.save()


# --- DOCX (python-docx) ---
import docx
from docx.shared import Pt


def gen_docx():
    path = os.path.join(HERE, "sample.docx")
    d = docx.Document()
    d.add_heading("OKF DOCX Fixture", level=0)
    d.add_paragraph("This is a paragraph in the DOCX fixture.")
    d.add_heading("Section One", level=1)
    d.add_paragraph("Content under section one.")
    t = d.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "Name"
    t.cell(0, 1).text = "Value"
    t.cell(1, 0).text = "alpha"
    t.cell(1, 1).text = "1"
    d.save(path)


# --- XLSX (openpyxl) ---
from openpyxl import Workbook


def gen_xlsx():
    path = os.path.join(HERE, "sample.xlsx")
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Item", "Qty"])
    ws.append(["apple", 3])
    ws.append(["banana", 5])
    wb.save(path)


# --- PPTX (python-pptx) ---
from pptx import Presentation


def gen_pptx():
    path = os.path.join(HERE, "sample.pptx")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "OKF PPTX Fixture"
    slide.placeholders[1].text = "First slide body content."
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Second Slide"
    prs.save(path)


# --- HTML ---
def gen_html():
    path = os.path.join(HERE, "sample.html")
    html = """<!DOCTYPE html>
<html>
<head><title>OKF HTML Fixture</title></head>
<body>
<h1>Main Heading</h1>
<p>A paragraph with a <a href="https://example.com">link</a>.</p>
<ul><li>item one</li><li>item two</li></ul>
</body>
</html>
"""
    with open(path, "w", encoding="utf-8") as f:
        f.write(html)


# --- CSV ---
def gen_csv():
    path = os.path.join(HERE, "sample.csv")
    with open(path, "w", encoding="utf-8") as f:
        f.write("name,value\nfoo,1\nbar,2\n")


# --- TXT ---
def gen_txt():
    path = os.path.join(HERE, "sample.txt")
    with open(path, "w", encoding="utf-8") as f:
        f.write("OKF TXT Fixture\nPlain text content line one.\n")


# --- Blank PDF (no text layer, for scanned-PDF error tests) ---
def gen_blank_pdf():
    path = os.path.join(HERE, "blank.pdf")
    c = canvas.Canvas(path, pagesize=letter)
    c.drawImage  # no-op marker to keep import
    c.showPage()  # page with zero text
    c.save()


if __name__ == "__main__":
    gen_pdf()
    gen_docx()
    gen_xlsx()
    gen_pptx()
    gen_html()
    gen_csv()
    gen_txt()
    gen_blank_pdf()
    print("fixtures generated in", HERE)
